#!/usr/bin/env python3
"""
Appodeal Pulse — PPTX Export

Converts SLIDES JSON into a beautiful PowerPoint presentation
styled with the BidMachine palette (dark theme, colored accents).
Archives presentations by date.
"""

import json
import os
import re
import sys
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import config

# ── BidMachine palette ────────────────────────────────────────

BG_COLOR = RGBColor(0x08, 0x08, 0x0C)       # #08080C
BG_CARD = RGBColor(0x14, 0x14, 0x1A)         # Card background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xE2, 0xE2, 0xEA)
DIM = RGBColor(0x8F, 0x95, 0xA3)
DARK_DIM = RGBColor(0x5C, 0x62, 0x70)


def hex_to_rgb(hex_str):
    """Convert #RRGGBB to RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ── Slide dimensions (16:9) ──────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_COLOR):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text) if text is not None else ""
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_badge(slide, text, accent_hex, top=Inches(0.8)):
    """Add a colored badge/tag at the top of a slide."""
    accent = hex_to_rgb(accent_hex)
    left = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(text) if text is not None else "".upper()
    p.font.size = Pt(10)
    p.font.color.rgb = accent
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_title(slide, text, accent_hex, top=Inches(1.3)):
    """Add a styled title."""
    accent = hex_to_rgb(accent_hex)
    return add_text_box(
        slide, Inches(1.2), top, Inches(10.9), Inches(0.8),
        text, font_size=32, color=accent, bold=True,
        font_name="Calibri", alignment=PP_ALIGN.LEFT
    )


def add_accent_bar(slide, accent_hex, top=Inches(0.4)):
    """Add a thin colored accent bar at the top."""
    accent = hex_to_rgb(accent_hex)
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), top, SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()


# ── Slide renderers ───────────────────────────────────────────

def render_title_slide(prs, date_str):
    """Create the opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Logo text
    add_text_box(
        slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(0.6),
        "APPODEAL", font_size=18, color=DARK_DIM, bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )

    # PULSE
    add_text_box(
        slide, Inches(1.2), Inches(2.7), Inches(10.9), Inches(1.2),
        "PULSE", font_size=56, color=WHITE, bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )

    # Date
    add_text_box(
        slide, Inches(1.2), Inches(4.2), Inches(10.9), Inches(0.5),
        date_str, font_size=16, color=DIM, bold=False,
        alignment=PP_ALIGN.CENTER
    )

    # Accent bar
    accent = hex_to_rgb("#6366F1")
    bar = slide.shapes.add_shape(1, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def render_birthday(slide, d):
    """Render birthday slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#FFB72B")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '🎂')}  CELEBRATION", accent)
    add_title(slide, d.get("title", "Happy Birthday!"), accent)

    # Name (big)
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(0.7),
        d.get("name", ""), font_size=28, color=WHITE, bold=True,
        alignment=PP_ALIGN.LEFT
    )
    # Date
    add_text_box(
        slide, Inches(1.2), Inches(3.0), Inches(10.9), Inches(0.4),
        d.get("date", ""), font_size=14, color=DIM,
        alignment=PP_ALIGN.LEFT
    )
    # Message
    add_text_box(
        slide, Inches(1.2), Inches(3.6), Inches(10.9), Inches(1.0),
        d.get("message", ""), font_size=16, color=SOFT_WHITE,
        alignment=PP_ALIGN.LEFT
    )
    # Team note
    add_text_box(
        slide, Inches(1.2), Inches(5.0), Inches(10.9), Inches(0.5),
        d.get("teamNote", ""), font_size=14, color=hex_to_rgb(accent),
        italic=True, alignment=PP_ALIGN.LEFT
    )


def render_win(slide, d):
    """Render win/achievement slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#12DF58")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '📈')}  WIN", accent)
    add_title(slide, d.get("title", "Achievement"), accent)

    # Big stat
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(5), Inches(0.9),
        d.get("stat", ""), font_size=48, color=hex_to_rgb(accent), bold=True
    )
    # Stat unit
    add_text_box(
        slide, Inches(1.2), Inches(3.1), Inches(5), Inches(0.4),
        d.get("statUnit", ""), font_size=12, color=DIM
    )
    # Who
    add_text_box(
        slide, Inches(1.2), Inches(3.6), Inches(10.9), Inches(0.4),
        f"by {d.get('who', '')}", font_size=14, color=SOFT_WHITE
    )
    # Headline
    add_text_box(
        slide, Inches(1.2), Inches(4.2), Inches(10.9), Inches(0.5),
        d.get("headline", ""), font_size=18, color=WHITE, bold=True
    )
    # Description
    add_text_box(
        slide, Inches(1.2), Inches(4.8), Inches(10.9), Inches(1.0),
        d.get("description", ""), font_size=14, color=SOFT_WHITE
    )
    # Takeaway
    if d.get("takeaway"):
        add_text_box(
            slide, Inches(1.2), Inches(6.0), Inches(10.9), Inches(0.5),
            f'💬 "{d["takeaway"]}"', font_size=13, color=DIM, italic=True
        )


def render_clap(slide, d):
    """Render clap/recognition slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#5D2CAC")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '👏')}  CLAPS", accent)
    add_title(slide, d.get("title", "Recognition"), accent)

    # Count
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(3), Inches(0.6),
        d.get("count", ""), font_size=28, color=hex_to_rgb(accent), bold=True
    )
    # From
    add_text_box(
        slide, Inches(1.2), Inches(2.9), Inches(10.9), Inches(0.4),
        f"from {d.get('from', '')}", font_size=14, color=SOFT_WHITE
    )
    # To
    to_list = d.get("to", [])
    if to_list:
        to_str = "to: " + ", ".join(to_list)
        add_text_box(
            slide, Inches(1.2), Inches(3.4), Inches(10.9), Inches(0.4),
            to_str, font_size=14, color=WHITE, bold=True
        )
    # Reason quote
    add_text_box(
        slide, Inches(1.2), Inches(4.1), Inches(10.9), Inches(1.2),
        f'"{d.get("reason", "")}"', font_size=15, color=SOFT_WHITE, italic=True
    )
    # Values
    values = d.get("values", [])
    if values:
        add_text_box(
            slide, Inches(1.2), Inches(5.6), Inches(10.9), Inches(0.5),
            " · ".join(values), font_size=12, color=hex_to_rgb(accent), bold=True
        )


def render_newjoin(slide, d):
    """Render new team member slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#00FFD1")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '🚀')}  NEW TEAMMATE", accent)
    add_title(slide, d.get("title", "Welcome!"), accent)

    # Name
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(0.7),
        d.get("name", ""), font_size=28, color=WHITE, bold=True
    )
    # Role + Team
    role_team = f"{d.get('role', '')} · {d.get('team', '')}"
    add_text_box(
        slide, Inches(1.2), Inches(3.0), Inches(10.9), Inches(0.4),
        role_team, font_size=14, color=hex_to_rgb(accent), bold=True
    )
    # Quote
    if d.get("quote"):
        add_text_box(
            slide, Inches(1.2), Inches(3.7), Inches(10.9), Inches(1.5),
            f'"{d["quote"]}"', font_size=15, color=SOFT_WHITE, italic=True
        )
    # Welcome
    add_text_box(
        slide, Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.5),
        "👋 Welcome aboard! 🎉", font_size=18, color=DIM
    )


def render_event(slide, d):
    """Render event slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#E84039")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '📅')}  EVENT", accent)
    add_title(slide, d.get("title", "Upcoming Event"), accent)

    # Event name
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(0.6),
        d.get("eventName", ""), font_size=22, color=WHITE, bold=True
    )
    # Date
    add_text_box(
        slide, Inches(1.2), Inches(3.0), Inches(10.9), Inches(0.4),
        d.get("date", ""), font_size=14, color=hex_to_rgb(accent), bold=True
    )
    # Description
    add_text_box(
        slide, Inches(1.2), Inches(3.6), Inches(10.9), Inches(1.5),
        d.get("description", ""), font_size=15, color=SOFT_WHITE
    )
    # Organizer
    if d.get("organizer"):
        add_text_box(
            slide, Inches(1.2), Inches(5.5), Inches(10.9), Inches(0.4),
            f"Organized by {d['organizer']}", font_size=13, color=DIM
        )


def render_milestone(slide, d):
    """Render milestone slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#1667EF")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '🏆')}  MILESTONE", accent)
    add_title(slide, d.get("title", "Milestone"), accent)

    # Big stat
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(5), Inches(0.9),
        d.get("stat", ""), font_size=52, color=hex_to_rgb(accent), bold=True
    )
    # Stat label
    add_text_box(
        slide, Inches(1.2), Inches(3.2), Inches(5), Inches(0.4),
        d.get("statLabel", ""), font_size=11, color=DIM
    )
    # Headline
    add_text_box(
        slide, Inches(1.2), Inches(3.9), Inches(10.9), Inches(0.5),
        d.get("headline", ""), font_size=18, color=WHITE, bold=True
    )
    # Description
    add_text_box(
        slide, Inches(1.2), Inches(4.6), Inches(10.9), Inches(1.5),
        d.get("description", ""), font_size=15, color=SOFT_WHITE
    )


def render_reading(slide, d):
    """Render reading/articles slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#5D2CAC")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '📚')}  RECOMMENDED", accent)
    add_title(slide, d.get("title", "What We're Reading"), accent)

    articles = d.get("articles", [])
    y_pos = 2.3
    for i, article in enumerate(articles[:3]):
        # Number
        add_text_box(
            slide, Inches(1.2), Inches(y_pos), Inches(0.5), Inches(0.4),
            f"0{i+1}", font_size=22, color=hex_to_rgb(accent), bold=True
        )
        # Title
        add_text_box(
            slide, Inches(1.9), Inches(y_pos), Inches(10.2), Inches(0.4),
            article.get("title", ""), font_size=16, color=WHITE, bold=True
        )
        # Description
        add_text_box(
            slide, Inches(1.9), Inches(y_pos + 0.35), Inches(10.2), Inches(0.4),
            article.get("desc", ""), font_size=12, color=SOFT_WHITE
        )
        # Shared by
        shared_info = f"Shared by {article.get('sharedBy', '')}"
        if article.get("reactions"):
            shared_info += f" · {article['reactions']}"
        add_text_box(
            slide, Inches(1.9), Inches(y_pos + 0.7), Inches(10.2), Inches(0.3),
            shared_info, font_size=11, color=DIM
        )
        y_pos += 1.3

    # Channel
    if d.get("channel"):
        add_text_box(
            slide, Inches(1.2), Inches(6.5), Inches(10.9), Inches(0.3),
            d["channel"], font_size=11, color=DARK_DIM
        )


def render_officelife(slide, d):
    """Render office life slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#00FFD1")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '💬')}  OFFICE LIFE", accent)
    add_title(slide, d.get("title", "Office Life"), accent)

    # Headline
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(0.5),
        d.get("headline", ""), font_size=18, color=WHITE, bold=True
    )
    # Quote
    if d.get("quote"):
        add_text_box(
            slide, Inches(1.5), Inches(3.0), Inches(10.6), Inches(2.0),
            f'"{d["quote"]}"', font_size=15, color=SOFT_WHITE, italic=True
        )
    # Author
    add_text_box(
        slide, Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.4),
        f"— {d.get('author', '')}", font_size=14, color=DIM
    )
    # Reactions
    if d.get("reactions"):
        add_text_box(
            slide, Inches(1.2), Inches(5.7), Inches(10.9), Inches(0.4),
            d["reactions"], font_size=14, color=hex_to_rgb(accent), bold=True
        )
    # Channel
    if d.get("channel"):
        add_text_box(
            slide, Inches(1.2), Inches(6.3), Inches(10.9), Inches(0.3),
            d["channel"], font_size=11, color=DARK_DIM
        )


def render_celebration(slide, d):
    """Render celebration/weekly ritual slide."""
    set_slide_bg(slide)
    accent = d.get("accent", "#FFB72B")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '🎉')}  WEEKLY RITUAL", accent)
    add_title(slide, d.get("title", "Celebration"), accent)

    # Headline
    add_text_box(
        slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(0.5),
        d.get("headline", ""), font_size=18, color=WHITE, bold=True
    )

    # Prompts
    y_pos = 3.2
    for prompt in [d.get("prompt1"), d.get("prompt2")]:
        if prompt:
            add_text_box(
                slide, Inches(1.5), Inches(y_pos), Inches(10.3), Inches(0.6),
                prompt, font_size=14, color=SOFT_WHITE
            )
            y_pos += 0.8

    # Core value
    if d.get("coreValue"):
        add_text_box(
            slide, Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.8),
            d["coreValue"], font_size=15, color=SOFT_WHITE, italic=True
        )
    # Channel
    if d.get("channel"):
        add_text_box(
            slide, Inches(1.2), Inches(6.3), Inches(10.9), Inches(0.3),
            d["channel"], font_size=11, color=DARK_DIM
        )


def render_generic(slide, d):
    """Fallback renderer for unknown slide types."""
    set_slide_bg(slide)
    accent = d.get("accent", "#888888")
    add_accent_bar(slide, accent)
    add_badge(slide, f"{d.get('emoji', '📌')}  {d.get('type', 'INFO').upper()}", accent)
    add_title(slide, d.get("title", "Update"), accent)

    if d.get("headline"):
        add_text_box(
            slide, Inches(1.2), Inches(2.3), Inches(10.9), Inches(0.5),
            d["headline"], font_size=18, color=WHITE, bold=True
        )
    if d.get("description"):
        add_text_box(
            slide, Inches(1.2), Inches(3.0), Inches(10.9), Inches(2.0),
            d["description"], font_size=15, color=SOFT_WHITE
        )


# ── Renderer dispatch ─────────────────────────────────────────

RENDERERS = {
    "birthday":    render_birthday,
    "win":         render_win,
    "clap":        render_clap,
    "newjoin":     render_newjoin,
    "event":       render_event,
    "milestone":   render_milestone,
    "reading":     render_reading,
    "officelife":  render_officelife,
    "celebration": render_celebration,
}


# ── Main export function ──────────────────────────────────────

def export_pptx(slides, output_path=None):
    """Convert SLIDES JSON array to a PPTX presentation.

    Args:
        slides: List of slide dictionaries
        output_path: Where to save. If None, archives by date.

    Returns:
        Path to the saved PPTX file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    today = datetime.now()
    date_str = today.strftime("%A, %B %d, %Y")
    render_title_slide(prs, date_str)

    # Content slides
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide_type = slide_data.get("type", "")
        renderer = RENDERERS.get(slide_type, render_generic)
        renderer(slide, slide_data)

    # Determine output path
    if output_path is None:
        archive_dir = os.path.join(config.DATA_DIR, "archive")
        os.makedirs(archive_dir, exist_ok=True)
        filename = f"pulse_{today.strftime('%Y-%m-%d')}.pptx"
        output_path = os.path.join(archive_dir, filename)

    prs.save(output_path)
    return output_path


# ── Entry point ─────────────────────────────────────────────

if __name__ == "__main__":
    # Load slides from data/slides.json
    slides_path = os.path.join(config.DATA_DIR, "slides.json")
    if not os.path.exists(slides_path):
        print(f"ERROR: {slides_path} not found. Run generate.py first.")
        sys.exit(1)

    with open(slides_path, "r", encoding="utf-8") as f:
        slides = json.load(f)

    output = export_pptx(slides)
    print(f"  Exported {len(slides)} slides → {output}")
